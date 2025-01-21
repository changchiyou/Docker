import io
import json
import os
from urllib.parse import quote_plus

import requests
from pptx import Presentation
from pptx.util import Inches

dir_path = "template"

API_KEY = os.getenv("PEXELS_API_KEY")


def parse_response(response):
    slides = response.split("\n\n")
    slides_content = []
    for slide in slides:
        lines = slide.split("\n")
        title_line = lines[0]
        if ": " in title_line:
            title = title_line.split(": ", 1)[1]  # Extract the title after 'Slide X: '
        else:
            title = title_line
        content_lines = [
            line.strip("-")
            for line in lines[1:]
            if line != "Content:" and line.startswith("-")
        ]  # Skip line if it is 'Content:'
        content = "\n".join(content_lines)  # Join the lines to form the content
        # Extract the keyword from the line that starts with 'Keyword:'
        keyword_line = [line for line in lines if "Keyword:" or "Keywords:" in line][0]
        keyword = keyword_line.split(": ", 1)[1]
        slides_content.append({"title": title, "content": content, "keyword": keyword})
    return slides_content


def search_pexels_images(keyword):
    query = quote_plus(keyword.lower())
    print("Query:", query)  # Debug
    PEXELS_API_URL = f"https://api.pexels.com/v1/search?query={query}&per_page=1"
    print("URL:", PEXELS_API_URL)  # Debug
    headers = {"Authorization": API_KEY}
    response = requests.get(PEXELS_API_URL, headers=headers)
    print("Response Status Code:", response.status_code)  # Debug
    print("Response Content:", response.text)  # Debug
    data = json.loads(response.text)
    if "photos" in data:
        if len(data["photos"]) > 0:
            return data["photos"][0]["src"]["medium"]
    return None


def delete_first_two_slides(presentation):
    slide_ids = [1, 0]
    for slide_id in slide_ids:
        if slide_id < len(presentation.slides):
            xml_slides = presentation.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[slide_id])


def create_ppt(
    slides_content,
    template_choice,
    presentation_title,
    presenter_name,
    insert_image,
    filename,
):
    current_dir = os.path.abspath(os.path.dirname(__file__))
    app_dir = os.path.abspath(os.path.join(current_dir, os.pardir))
    template_path = os.path.join(app_dir, "template", f"{template_choice}.pptx")

    prs = Presentation(template_path)

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # Extract initial formatting from the placeholders
    title_formatting = extract_text_formatting(title)
    subtitle_formatting = extract_text_formatting(subtitle)
    # Set text and apply formatting
    title.text = presentation_title
    set_text_formatting(title, title_formatting)

    subtitle.text = f"Presented by {presenter_name}"
    set_text_formatting(subtitle, subtitle_formatting)

    # add content slides
    for slide_content in slides_content:
        slide = prs.slides.add_slide(content_slide_layout)

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                content_title_formatting = extract_text_formatting(placeholder)
                placeholder.text = slide_content["title"]
                set_text_formatting(placeholder, content_title_formatting)

            elif placeholder.placeholder_format.type == 7:  # Content
                content_body_formatting = extract_text_formatting(placeholder)
                placeholder.text = slide_content["content"]
                set_text_formatting(placeholder, content_body_formatting)

        if insert_image:
            # fetch image URL from Pixabay based on the slide's title
            image_url = search_pexels_images(slide_content['keyword'])
            print("Image URL:", image_url) #debug
            if image_url is not None:
                # download the image
                image_data = requests.get(image_url).content
                # load image into BytesIO object
                image_stream = io.BytesIO(image_data)
                # add the image at the specified position
                slide_width = Inches(20)
                slide_height = Inches(15)

                image_width = Inches(8)  # width of image
                image_height = Inches(5)  # height of image

                left = slide_width - image_width  # calculate left position
                top = slide_height - image_height - Inches(4)  # calculate top position

                slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)

    # add credits slide
    if insert_image:
        slide = prs.slides.add_slide(content_slide_layout)
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                credits_title_formatting = extract_text_formatting(placeholder)
                placeholder.text = "Credits"
                set_text_formatting(placeholder, credits_title_formatting)

            elif placeholder.placeholder_format.type == 7:  # Content
                credits_body_formatting = extract_text_formatting(placeholder)
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                set_text_formatting(placeholder, credits_body_formatting)

    # Delete the first two slides after all new slides have been added
    delete_first_two_slides(prs)

    # Save the presentation
    prs.save(os.path.join(app_dir, "generated", f"{filename}.pptx"))


def extract_text_formatting(placeholder):
    formatting = {"font_name": None, "font_size": None, "font_color": None}
    if placeholder.text_frame:
        paragraph = placeholder.text_frame.paragraphs[0]
        if paragraph.runs:
            run = paragraph.runs[0]
            formatting["font_name"] = run.font.name
            formatting["font_size"] = run.font.size
            if run.font.color.rgb:
                formatting["font_color"] = run.font.color.rgb
    return formatting


def set_text_formatting(placeholder, formatting):
    for paragraph in placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            if formatting["font_name"]:
                run.font.name = formatting["font_name"]
            if formatting["font_size"]:
                run.font.size = formatting["font_size"]
            if formatting["font_color"]:
                run.font.color.rgb = formatting["font_color"]
